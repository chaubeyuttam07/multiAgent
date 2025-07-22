from pptx import Presentation

def extract_text_from_pptx(file_path):
    """Extract text from a PPTX file."""
    try:
        prs = Presentation(file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + " "
        return text
    except Exception as e:
        print(f"Error reading PPTX: {e}")
        return ""